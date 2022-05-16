from pptx import Presentation
from pptx.util import Inches, Pt
import cv2
import os
import re

project_directory = r'C:\Users\Micha.Vardy\Desktop\projects\threats\NL'
movie_directory = f'{project_directory}\\movies'
image_directory =  f"{project_directory}\\images"
os.chdir(project_directory)

def extract_image(movie_file,file_name):
    success,image  = cv2.VideoCapture(movie_file).read()
    image_path = f"{image_directory}\\{file_name}.jpg"
    cv2.imwrite(image_path,image)
    return(image_path)

def extract_aspect_ratio(image_path):
    src = cv2.imread(image_path)
    return(src.shape[1] / src.shape[0])


prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

for index, mov in enumerate(os.listdir(movie_directory)):
    index = index + 1
    if re.search('\.wmv',mov):
        # add slide
        slide = prs.slides.add_slide(blank_slide_layout)
        
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = mov
        p.font.bold = True
        p.font.size=Pt(30)

        movie_file = f"{movie_directory}\\{mov}"
        left = top = Inches(2)
        poster_frame_image = extract_image(movie_file,mov)
        ar = 2.0272727272727273
        width = Inches(6)
        height = Inches(3)
        slide.shapes.add_movie(movie_file, left,top,width,height,poster_frame_image)




prs.save(f"{project_directory}\\movies.pptx")